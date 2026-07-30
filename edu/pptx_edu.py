import os
import logging
import subprocess
try:
    from pptx import Presentation
except Exception:
    Presentation = None
from db.db_operations import insert_data, delete_data
from backend.config import Config
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from utils.logging_util import LoggerSingleton
import asyncio
from edu.file_text_extract_timeout import await_file_text_extract

# from db.db_redis import job_manager, job_progress_manager
from db.db_job_managers import AsyncJobManager, AsyncJobProgress


# 로거 설정
logger = LoggerSingleton.get_logger(logger_name="edu.ppt", level=logging.INFO)

embedding_model = OpenAIEmbeddings(openai_api_key=Config.OPENAI_API_KEY)


def convert_ppt_to_pptx(ppt_file_path):
    """PPT 파일을 PPTX로 변환"""
    output_dir = os.path.dirname(ppt_file_path)
    pptx_file_path = ppt_file_path.replace(".ppt", ".pptx")  # 변환될 파일 예상 경로
    try:
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pptx",
                ppt_file_path,
                "--outdir",
                output_dir,
            ],
            check=True,
        )
        # 변환 후 파일이 존재하는지 확인
        if not os.path.exists(pptx_file_path):
            logger.error(f"변환된 파일이 존재하지 않습니다: {pptx_file_path}")
            raise FileNotFoundError(
                f"변환된 파일이 존재하지 않습니다: {pptx_file_path}"
            )

        return pptx_file_path
    except subprocess.CalledProcessError as e:
        logger.error(f"파일 변환 중 오류 발생: {e}")
        raise


def extract_text_from_ppt(file_path: str) -> str:
    if Presentation is None:
        raise ImportError("python-pptx is required to process PPT/PPTX files")
    presentation = Presentation(file_path)
    text_content = []

    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_content = [f"\n=== Slide {slide_num} ===\n"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

            # 표 처리
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append("| " + " | ".join(row_data) + " |")
                slide_content.append("\n".join(table_data))

        text_content.extend(slide_content)

    return "\n".join(text_content)


async def process_ppt(
    content: str,
    file_path: str,
    table_name: str,
    dbname: str,
    job_id: str,
    each_progress: float,
    job_manager: AsyncJobManager,
    job_progress_manager: AsyncJobProgress,
    subject: str = "",
    memo: str = "",
    personal_info_filter: str = "N",  # 개인정보 필터링 옵션 추가
):
    """
    PPT/PPTX 파일을 처리하여 텍스트를 추출하고 벡터화한 후 데이터베이스에 저장합니다.
    """
    try:
        # PPT 파일인 경우 PPTX로 변환
        if content.lower().endswith(".ppt"):
            file_path = await await_file_text_extract(
                asyncio.to_thread(convert_ppt_to_pptx, file_path),
                path=file_path,
                stage="ppt_convert",
                logger=logger,
            )
            if not os.path.exists(file_path):
                raise FileNotFoundError(
                    f"PPTX 변환 후 파일이 존재하지 않습니다: {file_path}"
                )
        else:
            file_path = file_path

        # PPT/PPTX 파일에서 텍스트 추출
        extracted_text = await await_file_text_extract(
            asyncio.to_thread(extract_text_from_ppt, file_path),
            path=file_path,
            stage="pptx",
            logger=logger,
        )

        if not extracted_text.strip():
            logger.warning(f"파일에서 텍스트를 추출할 수 없습니다: {file_path}")
            return {"status": "warning", "message": "텍스트를 추출할 수 없습니다."}

        # 텍스트 분할
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=Config.BASIC_CHUNK_SIZE, chunk_overlap=Config.BASIC_CHUNK_OVERLAP
        )
        chunks = text_splitter.split_text(extracted_text)
        total_chunks = len(chunks)
        # ✅ main.py에서 이미 전체 청크 수 기준으로 계산된 청크별 진행률을 그대로 사용
        chunk_progress = each_progress
        logger.info(
            f"[conent: {content}, total_chunk: {total_chunks}, 학습 시작 job_id:{job_id}, table_name: {table_name}]"
        )
        # 청크 처리 및 저장
        for idx, chunk in enumerate(chunks):
            status = await job_manager.get_job_status(job_id)
            if status == "cancel":
                logger.info(f"작업이 취소되었습니다: job_id={job_id}")
                return
            if personal_info_filter == "Y":
                logger.info(f"PPT 개별 청크 처리 중 개인정보 필터링 적용: job_id={job_id}, chunk={idx+1}")
                from utils.dlp_api import check_pii_content
                original_chunk = chunk  # 원본 청크 저장
                pii_result = check_pii_content(chunk)
                
                # 마스킹된 텍스트를 chunk 변수에 할당
                if pii_result["success"]:
                    chunk = pii_result["masked_text"]
                    

                else:
                    # PII 검사 실패 시 에러 로깅
                    logger.error(f"PII 검사 실패: {pii_result.get('error', 'Unknown error')}")
                    # 실패 시 원본 텍스트 유지
            chunk_num = f"{idx + 1}"
            chunk_with_metadata = chunk
            embedding = embedding_model.embed_query(chunk_with_metadata)

            # PostgreSQL vector 형식으로 변환 (대괄호[] 사용)
            embedding_array = f"[{','.join(map(str, embedding))}]"

            await insert_data(
                table=table_name,
                data={
                    "content": content,
                    "subject": subject or os.path.basename(file_path or "") or content,
                    "chunk_num": chunk_num,
                    "memo": memo,
                    "content_type": "file",
                    "text_data": chunk_with_metadata,
                    "embedding": embedding_array,
                },
                dbname=dbname,
            )
            # 진행률 업데이트
            current_progress = await job_progress_manager.get_job_progress(job_id)
            new_progress = round(min(current_progress + chunk_progress, 99.99), 2)
            await job_progress_manager.set_job_progress(job_id, new_progress)

            # logger.info(f"[Chunk_number: {idx + 1}]")

        logger.info(
            f"[conent: {content}, total_chunk: {total_chunks}, 학습 완료 job_id:{job_id}, table_name: {table_name}]"
        )
        return {
            "status": "success", 
            "message": f"{content} 처리 완료", 
            "chunks": total_chunks,
            "chunk_count": [total_chunks],
            "use_source": [content]
        }

    except Exception as e:
        logger.error(f"PPT 파일 처리 중 오류 발생: {e}", exc_info=True)

        raise RuntimeError(f"PPT 파일 처리 중 오류 발생: {e}")


async def calculate_pptx_chunks(file_path: str) -> int:
    """PPTX 파일의 정확한 청크 수를 계산합니다."""
    try:
        import asyncio
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        from backend.config import Config
        
        # ✅ extract_text_from_ppt 함수와 동일한 로직으로 텍스트 추출
        def extract_text_sync(file_path: str) -> str:
            presentation = Presentation(file_path)
            text_content = []

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = [f"\n=== Slide {slide_num} ===\n"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                    # 표 처리
                    if hasattr(shape, "has_table") and shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append("| " + " | ".join(row_data) + " |")
                        slide_content.append("\n".join(table_data))

                text_content.extend(slide_content)

            return "\n".join(text_content)
        
        # 비동기로 텍스트 추출 실행
        extracted_text = await asyncio.to_thread(extract_text_sync, file_path)
        
        if not extracted_text.strip():
            logger.warning(f"PPTX 파일에서 텍스트를 추출할 수 없습니다: {file_path}")
            return 1  # 기본값
        
        # ✅ 동일한 청크 분할 설정 사용
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=Config.BASIC_CHUNK_SIZE, 
            chunk_overlap=Config.BASIC_CHUNK_OVERLAP
        )
        chunks = text_splitter.split_text(extracted_text)
        actual_chunks = len(chunks)
        
        logger.info(f"PPTX 정확한 청크 수 계산: {file_path}, 실제 청크: {actual_chunks}개")
        return actual_chunks
        
    except Exception as e:
        logger.error(f"PPTX 청크 수 계산 실패: {file_path}, 오류: {e}")
        return 5  # 기본값
